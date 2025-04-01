import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
import docx  # Library to handle .docx files
from pptx import Presentation  # Library to handle .pptx files
import csv  # Built-in library to handle .csv files

# Hardcode your Google API key here (replace with your actual key)
GOOGLE_API_KEY = "AIzaSyBSexVrceATa6AnjuTgVPefOzUxyi2cvEM"

def get_document_text(docs):
    """Extract text from PDF, DOCX, TXT, PPTX, and CSV files"""
    text = ""
    for doc in docs:
        try:
            if doc.name.endswith('.pdf'):
                pdf_reader = PdfReader(doc)
                for page in pdf_reader.pages:
                    extracted_text = page.extract_text()
                    if extracted_text:
                        text += extracted_text
            elif doc.name.endswith('.docx'):
                docx_file = docx.Document(doc)
                for para in docx_file.paragraphs:
                    if para.text:
                        text += para.text + "\n"
            elif doc.name.endswith('.txt'):
                text += doc.read().decode('utf-8')  # Read text file content
            elif doc.name.endswith('.pptx'):
                prs = Presentation(doc)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + "\n"
            elif doc.name.endswith('.csv'):
                csv_content = doc.read().decode('utf-8')
                csv_reader = csv.reader(csv_content.splitlines(), delimiter=',')
                for row in csv_reader:
                    text += " ".join(row) + "\n"
        except Exception as e:
            st.error(f"Error processing file {doc.name}: {str(e)}")
    return text

def get_text_chunks(text):
    """Split text into chunks"""
    try:
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=10000,
            chunk_overlap=1000,
            length_function=len
        )
        chunks = text_splitter.split_text(text)
        return chunks if chunks else []
    except Exception as e:
        st.error(f"Error splitting text: {str(e)}")
        return []

def get_vector_store(text_chunks, api_key):
    """Create and save vector store"""
    try:
        embeddings = GoogleGenerativeAIEmbeddings(
            model="models/embedding-001",
            google_api_key=api_key
        )
        # Test embedding with a small sample
        embeddings.embed_query("test")
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
        vector_store.save_local("faiss_index")
        return True
    except Exception as e:
        st.error(f"Error creating vector store: {str(e)}")
        return False

def get_conversational_chain(api_key):
    """Create conversational chain"""
    prompt_template = """
    Answer the question as detailed as possible from the provided context.
    If the answer is not in the provided context, say "The answer is not available in the context".
    Do not provide incorrect information.
    
    Context: {context}
    Question: {question}
    
    Answer:
    """

    try:
        model = ChatGoogleGenerativeAI(
            model="gemini-1.5-flash",
            google_api_key=api_key,
            temperature=0.3
        )
        prompt = PromptTemplate(
            template=prompt_template,
            input_variables=["context", "question"]
        )
        chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
        return chain
    except Exception as e:
        st.error(f"Error creating conversational chain: {str(e)}")
        return None

def user_input(user_question, api_key):
    """Process user question and generate response"""
    try:
        if not os.path.exists("faiss_index"):
            st.error("No FAISS index found. Please upload and process files first.")
            return

        embeddings = GoogleGenerativeAIEmbeddings(
            model="models/embedding-001",
            google_api_key=api_key
        )
        new_db = FAISS.load_local(
            "faiss_index",
            embeddings,
            allow_dangerous_deserialization=True
        )
        docs = new_db.similarity_search(user_question, k=4)

        chain = get_conversational_chain(api_key)
        if chain is None:
            return

        response = chain(
            {"input_documents": docs, "question": user_question},
            return_only_outputs=True
        )
        
        st.write("Reply:", response["output_text"])
    except Exception as e:
        st.error(f"Error processing question: {str(e)}")

def main():
    """Main application function"""
    st.set_page_config(page_title="Chat with Documents", layout="wide")
    
    api_key = GOOGLE_API_KEY

    if not api_key or api_key == "your-actual-api-key-here":
        st.error("Please set a valid GOOGLE_API_KEY in the code.")
        st.info("Replace 'your-actual-api-key-here' with your actual Google API key at the top of the script.")
        return

    try:
        genai.configure(api_key=api_key)
        # Test API key with a simple call
        genai.get_model("models/gemini-1.5-flash")
    except Exception as e:
        st.error(f"Error validating API key: {str(e)}")
        st.info("Please verify your API key in Google Cloud Console and ensure Generative Language API is enabled.")
        return

    st.header("Chat with PDF, Documents, Text, PPTX & CSV files using Gemini💁")

    if 'processed' not in st.session_state:
        st.session_state.processed = False

    user_question = st.text_input("Ask a Question from the Uploaded Files", key="question_input")
    if user_question and st.session_state.processed:
        user_input(user_question, api_key)
    elif user_question and not st.session_state.processed:
        st.warning("Please upload and process files first")

    with st.sidebar:
        st.title("Menu")
        uploaded_docs = st.file_uploader(
            "Upload your Files (PDF, DOCX, TXT, PPTX, CSV)",
            accept_multiple_files=True,
            type=['pdf', 'docx', 'txt', 'pptx', 'csv'],
            key="doc_uploader"
        )
        
        if st.button("Submit & Process", key="process_button"):
            if not uploaded_docs:
                st.warning("Please upload at least one file")
            else:
                with st.spinner("Processing..."):
                    raw_text = get_document_text(uploaded_docs)
                    if not raw_text:
                        st.error("No text could be extracted from the files")
                    else:
                        text_chunks = get_text_chunks(raw_text)
                        if not text_chunks:
                            st.error("Failed to split text into chunks")
                        elif get_vector_store(text_chunks, api_key):
                            st.session_state.processed = True
                            st.success("Processing complete! You can now ask questions.")
                        else:
                            st.error("Failed to create vector store")

if __name__ == "__main__":
    main()
